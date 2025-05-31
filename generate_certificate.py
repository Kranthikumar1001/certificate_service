from pptx import Presentation
import os
import uuid
import subprocess

def download_template(template_url):
    import requests
    local_path = f"/tmp/{uuid.uuid4()}.pptx"
    r = requests.get(template_url)
    with open(local_path, 'wb') as f:
        f.write(r.content)
    return local_path

def generate_certificate(name, address, template_url):
    template_path = download_template(template_url)
    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "{{Name}}" in run.text:
                            run.text = run.text.replace("{{Name}}", name)
                        if "{{Address}}" in run.text:
                            run.text = run.text.replace("{{Address}}", address)

    output_pptx = f"/tmp/{uuid.uuid4()}.pptx"
    prs.save(output_pptx)

    # Convert PPTX to PDF using LibreOffice
    subprocess.run([
        'libreoffice', '--headless', '--convert-to', 'pdf', '--outdir',
        '/tmp', output_pptx
    ], check=True)

    output_pdf = output_pptx.replace(".pptx", ".pdf")
    return output_pptx, output_pdf
